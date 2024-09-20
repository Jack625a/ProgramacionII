from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Portada
slide_layout = prs.slide_layouts[0]  # Use the title slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Una mirada reflexiva a la educación superior y la docencia universitaria"
subtitle.text = "Juan Carlos Arroyo Mendizábal Ph.D."

# Slide 2: Índice
slide_layout = prs.slide_layouts[1]  # Use the title and content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Índice"
content = slide.placeholders[1].text = "\n".join([
    "1. Historia de la Universidad Latinoamericana",
    "2. Influencia Europea",
    "3. Primera Reforma de la Educación Superior",
    "4. Segunda Reforma de la Educación Superior",
    "5. Tercera Reforma de la Educación Superior",
    "6. Visión Actual de la Universidad Boliviana",
    "7. Universidades Públicas y Privadas",
    "8. Evaluación y Acreditación",
    "9. Filosofía Humanista y Formación Continua",
    "10. Desafíos en el Siglo XXI",
    "11. Visión de la UNESCO",
    "12. Futuros de la Educación Superior (UNESCO - IESALC)",
    "13. Conclusión y Reflexiones"
])

# Create placeholder slides for the rest
sections = [
    "Historia de la Universidad Latinoamericana",
    "Influencia Europea",
    "Primera Reforma de la Educación Superior",
    "Segunda Reforma de la Educación Superior",
    "Tercera Reforma de la Educación Superior",
    "Visión Actual de la Universidad Boliviana",
    "Universidades Públicas y Privadas",
    "Evaluación y Acreditación",
    "Filosofía Humanista y Formación Continua",
    "Desafíos en el Siglo XXI",
    "Visión de la UNESCO",
    "Futuros de la Educación Superior (UNESCO - IESALC)",
    "Conclusión y Reflexiones"
]

for section in sections:
    slide_layout = prs.slide_layouts[1]  # Use the title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = section
    content = slide.placeholders[1].text = "Contenido en desarrollo..."

# Save the presentation
prs.save("/mnt/data/Presentacion_Educacion_Superior.pptx")
3